from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLOR_AZURE = RGBColor(0, 120, 212)
COLOR_DARK = RGBColor(32, 32, 32)
COLOR_MUTED = RGBColor(90, 90, 90)
COLOR_LIGHT = RGBColor(242, 246, 252)
COLOR_BORDER = RGBColor(210, 220, 230)
COLOR_ACCENT = RGBColor(255, 140, 0)
COLOR_WHITE = RGBColor(255, 255, 255)

FONT_JP = "Meiryo"

EVENT_TITLE = (
    "Action required: Transition Azure Key Vault access policies to Azure RBAC "
    "or configure Azure Key Vault to explicitly use access policies"
)

TRACKING_ID = "RN3T-JRG"
SERVICE = "Azure Key Vault"
EVENT_TYPE = "正常性の勧告"
STATUS = "Active"
EVENT_LEVEL = "Informational"
EVENT_TAG = "Action Recommended"
START_TIME = "2026-02-09 09:00 JST"
LAST_UPDATE = "2026-02-10 04:28 JST"
RETIRE_DATE = "2027-02-27"

REGION_SUMMARY = (
    "主要な Azure パブリックリージョンが広く対象"
    "（East US / West US 2 / West Europe / North Europe / Japan East / Japan West / "
    "Southeast Asia / Australia East ほか多数）"
)

SOURCE_TEXT = "Source: Azure Service Health / Tracking ID: RN3T-JRG"


def set_run_style(run, size: int, bold: bool = False, color: RGBColor = COLOR_DARK):
    run.font.name = FONT_JP
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_cell_text(cell, text: str, size: int = 16, bold: bool = False, color: RGBColor = COLOR_DARK):
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run_style(run, size=size, bold=bold, color=color)


def add_footer(slide):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.2), Inches(0.25))
    p = tx.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = SOURCE_TEXT
    set_run_style(run, size=9, color=COLOR_MUTED)
    p.alignment = PP_ALIGN.RIGHT


def add_header(slide, title: str, subtitle: str | None = None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.65)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_AZURE
    band.line.color.rgb = COLOR_AZURE

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(10.8), Inches(0.32))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_style(run, size=24, bold=True, color=COLOR_WHITE)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(10.7), Inches(0.15), Inches(2.1), Inches(0.25))
        p2 = sub_box.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        set_run_style(run2, size=10, color=COLOR_WHITE)
        p2.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_WHITE

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.8)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_AZURE
    band.line.color.rgb = COLOR_AZURE

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11.8), Inches(1.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Azure Service Health 公開情報の紹介"
    set_run_style(run, size=28, bold=True, color=COLOR_DARK)

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Azure Key Vault access policies → Azure RBAC への移行勧告"
    set_run_style(run2, size=22, color=COLOR_AZURE)

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(3.0),
        Inches(5.9),
        Inches(2.6),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT
    card.line.color.rgb = COLOR_BORDER

    info = [
        f"Tracking ID: {TRACKING_ID}",
        f"Service: {SERVICE}",
        f"Event Type: {EVENT_TYPE}",
        f"Status: {STATUS}",
        f"Tag: {EVENT_TAG}",
        f"Start: {START_TIME}",
        f"Last update: {LAST_UPDATE}",
    ]

    text_box = slide.shapes.add_textbox(Inches(0.95), Inches(3.28), Inches(5.3), Inches(2.05))
    tf2 = text_box.text_frame
    tf2.clear()
    for i, line in enumerate(info):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        run = p.add_run()
        run.text = line
        set_run_style(run, size=16, color=COLOR_DARK)

    callout = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(7.0),
        Inches(3.15),
        Inches(5.0),
        Inches(2.3),
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = COLOR_WHITE
    callout.line.color.rgb = COLOR_ACCENT
    callout.line.width = Pt(2)

    t = slide.shapes.add_textbox(Inches(7.3), Inches(3.45), Inches(4.4), Inches(1.7))
    tf3 = t.text_frame
    p = tf3.paragraphs[0]
    run = p.add_run()
    run.text = "ポイント"
    set_run_style(run, size=18, bold=True, color=COLOR_ACCENT)

    bullets = [
        "2027-02-27 に旧 API が廃止",
        "新規 Key Vault は Azure RBAC が既定",
        "IaC / 自動化の前提見直しが必要",
    ]
    for b in bullets:
        p = tf3.add_paragraph()
        run = p.add_run()
        run.text = f"• {b}"
        set_run_style(run, size=16, color=COLOR_DARK)

    add_footer(slide)


def add_bullets_slide(prs: Presentation, title: str, bullets: list[tuple[int, str]], subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_WHITE

    add_header(slide, title, subtitle)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.7), Inches(5.7))
    tf = body.text_frame
    tf.clear()

    for i, (level, text) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(8)

        run = p.add_run()
        prefix = "• " if level == 0 else "– "
        run.text = f"{prefix}{text}"
        set_run_style(
            run,
            size=20 if level == 0 else 16,
            bold=(level == 0),
            color=COLOR_DARK if level == 0 else COLOR_MUTED,
        )

    add_footer(slide)


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]], subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_WHITE

    add_header(slide, title, subtitle)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.7), Inches(1.2), Inches(12.0), Inches(5.2)
    )
    table = table_shape.table

    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(9.0)

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_AZURE
        set_cell_text(cell, header, size=17, bold=True, color=COLOR_WHITE)

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_LIGHT if r % 2 == 1 else COLOR_WHITE
            set_cell_text(cell, value, size=15)

    add_footer(slide)


def build_presentation(output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs)

    add_bullets_slide(
        prs,
        "要点サマリー",
        [
            (0, f"{RETIRE_DATE} に Azure Key Vault API version 2026-02-01 より前の API が廃止される"),
            (0, "API version 2026-02-01 以降は、新規作成される Key Vault の既定アクセス制御モデルが Azure RBAC になる"),
            (0, "既存の Key Vault は現在のアクセス制御モデルを継続する"),
            (0, "Azure portal の挙動は変更されない"),
            (0, "access policies 前提の自動化を継続する場合は、明示的な設定が必要"),
            (0, "対応しないと HTTP 403 やロール不足による処理失敗の可能性がある"),
        ],
        subtitle="Summary",
    )

    add_bullets_slide(
        prs,
        "変更内容の詳細",
        [
            (0, "2026年2月リリースの API version 2026-02-01 で、セキュリティ上の重要な更新が入る"),
            (1, "新規作成されるすべての Key Vault で Azure RBAC が既定となる"),
            (1, "既存 Vault は今のアクセス制御モデルを維持する"),
            (0, "legacy access policies を使っている場合は、Azure RBAC への移行が推奨される"),
            (0, "access policies を継続したい場合は、以下で明示設定が必要"),
            (1, "CLI / PowerShell / REST API / ARM / Bicep / Terraform"),
            (0, "放置すると新規 Vault 作成時に RBAC 既定となり、権限不足の不具合が顕在化しうる"),
        ],
        subtitle="Details",
    )

    add_table_slide(
        prs,
        "影響範囲",
        headers=["観点", "内容"],
        rows=[
            ["対象サービス", "Azure Key Vault"],
            ["対象リージョン", REGION_SUMMARY],
            ["影響を受けやすいケース", "CLI / PowerShell / REST API / ARM / Bicep / Terraform で Key Vault を新規作成・再作成しているケース"],
            ["特に注意したい運用", "CI/CD、IaC、スクリプト化されたシークレット登録、環境再構築、自動プロビジョニング"],
            ["影響が比較的小さいケース", "既存 Vault の継続利用中心、Portal 手動操作中心、すでに Azure RBAC へ移行済み"],
            ["主な失敗モード", "HTTP 403、ロール不足、シークレット取得失敗、デプロイ失敗、運用ジョブ失敗"],
        ],
        subtitle="Impact",
    )

    add_bullets_slide(
        prs,
        "必要アクションと期限",
        [
            (0, "原則対応: 新規・既存 Vault を Azure RBAC へ移行する"),
            (0, "代替対応: access policies を継続したい場合は、新規 Vault 作成時に明示的に設定する"),
            (0, f"移行期限: {RETIRE_DATE} までに API version 2026-02-01 へ移行する"),
            (0, "実施ステップ"),
            (1, "Key Vault 利用実態の棚卸し"),
            (1, "access policies 利用有無の確認"),
            (1, "ARM / Bicep / Terraform / CLI / PowerShell のテンプレート確認"),
            (1, "検証環境で新規 Vault 作成と権限動作を確認"),
        ],
        subtitle="Required Action",
    )

    add_bullets_slide(
        prs,
        "自社としての対応方針",
        [
            (0, "本件は緊急障害対応ではないが、中期的には確実に対応が必要な設計変更案件"),
            (0, "特に Key Vault を自動構築しているチームでは優先度は中〜高"),
            (0, "短期"),
            (1, "利用システム / サブスクリプション / テンプレートの棚卸し"),
            (1, "access policies 前提のフローを洗い出し"),
            (0, "中期"),
            (1, "Azure RBAC への移行可否を判断"),
            (1, "必要ロール設計とテンプレート更新"),
            (0, "推奨メッセージ"),
            (1, "『既存 Vault はすぐ止まらないが、将来の新規構築でハマる前に整理する』"),
        ],
        subtitle="Recommendation",
    )

    add_table_slide(
        prs,
        "イベント情報（参考）",
        headers=["項目", "値"],
        rows=[
            ["Title", EVENT_TITLE],
            ["Tracking ID", TRACKING_ID],
            ["Event Type", EVENT_TYPE],
            ["Status", STATUS],
            ["Event level", EVENT_LEVEL],
            ["Event tag", EVENT_TAG],
            ["Service", SERVICE],
            ["Start time", START_TIME],
            ["Last update", LAST_UPDATE],
            ["Retirement date", RETIRE_DATE],
        ],
        subtitle="Appendix",
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def parse_args():
    parser = argparse.ArgumentParser(
        description="Generate a PowerPoint deck for Azure Service Health event RN3T-JRG."
    )
    parser.add_argument(
        "-o",
        "--output",
        default=r"c:\Users\hmatsumoto\OneDrive - Microsoft\03_git\security\logicapps\Azure_Service_Health_KeyVault_RN3T-JRG.pptx",
        help="Output PowerPoint path (.pptx)",
    )
    return parser.parse_args()


def main():
    args = parse_args()
    output_path = Path(args.output).resolve()
    created = build_presentation(output_path)
    print(f"PowerPoint created: {created}")


if __name__ == "__main__":
    main()